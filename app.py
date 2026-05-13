import streamlit as st
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image
import io
import json
import re

st.set_page_config(page_title="원본 재현 PPT 변환기", layout="wide")
st.title("🎨 원본 재현형 이미지-PPT 변환기")

# 1. API 키 설정 (Secrets에서 자동 로드)
if "GEMINI_API_KEY" in st.secrets:
    api_key = st.secrets["GEMINI_API_KEY"]
else:
    api_key = st.sidebar.text_input("Gemini API Key", type="password")

if api_key:
    genai.configure(api_key=api_key)
    # 모델 설정 (가장 안정적인 모델 명칭 사용)
    model = genai.GenerativeModel('gemini-1.5-flash')

    uploaded_files = st.file_uploader("슬라이드 이미지들을 업로드하세요", accept_multiple_files=True, type=['jpg', 'png', 'jpeg'])

    if uploaded_files and st.button("🚀 원본 재현 PPT 생성 시작"):
        prs = Presentation()
        # 슬라이드 크기 설정 (표준 16:9 비율)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        with st.spinner('이미지 분석 및 PPT 레이어 작업 중...'):
            for uploaded_file in uploaded_files:
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                
                # 2. 원본 이미지 읽기 및 슬라이드 배경으로 삽입
                img_data = uploaded_file.read()
                img_stream = io.BytesIO(img_data)
                slide.shapes.add_picture(img_stream, 0, 0, width=prs.slide_width, height=prs.slide_height)
                
                # 3. AI에게 텍스트 및 좌표 데이터 요청
                img = Image.open(io.BytesIO(img_data))
                prompt = """
                Extract all text with coordinates. Return ONLY a JSON list like this:
                [{"text": "Sample text", "x": 10, "y": 20, "w": 30, "h": 5}]
                The values x, y, w, h are percentages (0-100) of the image size.
                Return only the raw JSON list, no markdown tags like ```json.
                """
                
                try:
                    response = model.generate_content([prompt, img])
                    raw_text = response.text.strip()
                    
                    # [SyntaxError 해결 포인트] 정규표현식을 한 줄로 정확히 작성
                    clean_json = re.sub(r'
```(?:json)?|```', '', raw_text).strip()
                    
                    # 만약 JSON 형식이 아닐 경우 리스트 대괄호만 추출
                    match = re.search(r'\[.*\]', clean_json, re.DOTALL)
                    if match:
                        clean_json = match.group()
                        
                    text_blocks = json.loads(clean_json)
                    
                    # 4. 분석된 좌표에 텍스트 박스 얹기
                    for block in text_blocks:
                        # 좌표 데이터를 PPT 단위로 환산
                        bx = float(block.get('x', 0))
                        by = float(block.get('y', 0))
                        bw = float(block.get('w', 10))
                        bh = float(block.get('h', 5))
                        
                        left = prs.slide_width * (bx / 100)
                        top = prs.slide_height * (by / 100)
                        width = prs.slide_width * (bw / 100)
                        height = prs.slide_height * (bh / 100)
                        
                        txBox = slide.shapes.add_textbox(left, top, width, height)
                        tf = txBox.text_frame
                        tf.word_wrap = True
                        p = tf.add_paragraph()
                        p.text = block.get('text', '')
                        # 글꼴 크기 및 가독성 설정
                        p.font.size = Pt(14)
                        p.font.bold = True
                        
                except Exception as e:
                    st.warning(f"{uploaded_file.name}: 텍스트 추출에 실패했습니다. (이미지만 삽입됨) 사유: {e}")

            # 5. 최종 PPT 파일 생성 및 다운로드 제공
            ppt_out = io.BytesIO()
            prs.save(ppt_out)
            st.success("✅ PPT 생성이 완료되었습니다!")
            st.download_button(
                label="📥 완성된 PPT 다운로드",
                data=ppt_out.getvalue(),
                file_name="reproduced_slide.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
else:
    st.info("사이드바에 API 키를 입력하거나 Secrets에 GEMINI_API_KEY를 등록해주세요.")
